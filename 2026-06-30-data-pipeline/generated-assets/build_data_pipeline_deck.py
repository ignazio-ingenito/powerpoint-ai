from pathlib import Path
import shutil
import zipfile

from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Cm, Pt


ROOT = Path(__file__).resolve().parents[1]
REPO = ROOT.parent
TEMPLATE = REPO / "docs" / "template.pptx"
MEDIA_LOGO = Path("/tmp/pptx-media/ppt/media/image15.png")
OUT = ROOT / "2026_TXT_001 - Data Pipeline Blueprint - Scenari TO BE v0.3.pptx"
ASSET_LOGO = ROOT / "generated-assets" / "txt-novigo-logo-from-template.png"

W = 13.333333
H = 7.5

COLORS = {
    "teal": RGBColor(16, 164, 168),
    "teal_dark": RGBColor(0, 122, 130),
    "blue": RGBColor(43, 137, 203),
    "navy": RGBColor(31, 65, 95),
    "text": RGBColor(36, 42, 48),
    "muted": RGBColor(103, 112, 120),
    "line": RGBColor(149, 207, 211),
    "light": RGBColor(238, 249, 250),
    "very_light": RGBColor(247, 251, 252),
    "white": RGBColor(255, 255, 255),
    "gray": RGBColor(229, 235, 238),
    "green": RGBColor(70, 151, 103),
    "amber": RGBColor(215, 155, 63),
    "red": RGBColor(178, 88, 82),
}


def ensure_logo():
    if MEDIA_LOGO.exists():
        shutil.copyfile(MEDIA_LOGO, ASSET_LOGO)
    elif not ASSET_LOGO.exists():
        with zipfile.ZipFile(TEMPLATE) as zf:
            with zf.open("ppt/media/image15.png") as src:
                ASSET_LOGO.write_bytes(src.read())


def clear_template_slides(prs):
    slide_id_list = prs.slides._sldIdLst
    for slide_id in list(slide_id_list):
        r_id = slide_id.rId
        prs.part.drop_rel(r_id)
        slide_id_list.remove(slide_id)


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_textbox(slide, x, y, w, h, text, size=14, color="text", bold=False,
                align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Cm(0.05)
    tf.margin_right = Cm(0.05)
    tf.margin_top = Cm(0.02)
    tf.margin_bottom = Cm(0.02)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    for run in p.runs:
        run.font.name = "Aptos"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = COLORS[color]
    return box


def add_multiline(slide, x, y, w, h, lines, size=11, color="text",
                  bullet=False, gap=0):
    box = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Cm(0.08)
    tf.margin_right = Cm(0.08)
    tf.margin_top = Cm(0.04)
    tf.margin_bottom = Cm(0.04)
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.space_after = Pt(gap)
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = COLORS[color]
        if bullet:
            p.text = f"- {line}"
    return box


def add_header(slide, section, num):
    slide.shapes.add_picture(str(ASSET_LOGO), Cm(0.35), Cm(0.23), width=Cm(2.1))
    add_textbox(slide, 2.55, 0.28, 10.0, 0.35, section, size=8, color="muted")
    add_textbox(slide, 32.55, 0.25, 0.45, 0.35, str(num), size=8, color="muted", align=PP_ALIGN.RIGHT)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(0.35), Cm(0.82), Cm(1.25), Cm(0.035))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["teal"]
    line.line.fill.background()


def add_title(slide, title, subtitle=None):
    add_textbox(slide, 1.0, 1.15, 31.2, 0.75, title, size=21, color="navy", bold=True)
    if subtitle:
        add_textbox(slide, 1.02, 1.9, 30.6, 0.55, subtitle, size=10.5, color="muted")


def rounded_rect(slide, x, y, w, h, fill="white", line="line", radius=True):
    shp = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Cm(x), Cm(y), Cm(w), Cm(h)
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = COLORS[fill]
    shp.line.color.rgb = COLORS[line]
    shp.line.width = Pt(1.0)
    return shp


def card(slide, x, y, w, h, title, body, accent="teal", body_size=9.5):
    rounded_rect(slide, x, y, w, h, fill="white", line="line")
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(x), Cm(y), Cm(w), Cm(0.42))
    band.fill.solid()
    band.fill.fore_color.rgb = COLORS[accent]
    band.line.fill.background()
    add_textbox(slide, x + 0.18, y + 0.08, w - 0.35, 0.3, title, size=8.5, color="white", bold=True)
    if isinstance(body, list):
        add_multiline(slide, x + 0.18, y + 0.62, w - 0.35, h - 0.75, body, size=body_size, color="text", bullet=True)
    else:
        add_multiline(slide, x + 0.18, y + 0.62, w - 0.35, h - 0.75, [body], size=body_size, color="text")


def pill(slide, x, y, w, text, fill="light", line="teal", size=8.5):
    shp = rounded_rect(slide, x, y, w, 0.48, fill=fill, line=line)
    add_textbox(slide, x + 0.1, y + 0.12, w - 0.2, 0.25, text, size=size, color="navy", bold=True, align=PP_ALIGN.CENTER)
    return shp


def connector(slide, x1, y1, x2, y2, color="line"):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(x1), Cm(y1), Cm(x2), Cm(y2))
    c.line.color.rgb = COLORS[color]
    c.line.width = Pt(1.2)
    return c


def cover(prs):
    s = blank_slide(prs)
    s.shapes.add_picture(str(ASSET_LOGO), Cm(0.65), Cm(0.45), width=Cm(2.5))
    deco = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ARC, Cm(24.0), Cm(0.8), Cm(7.8), Cm(5.8))
    deco.line.color.rgb = COLORS["teal"]
    deco.line.width = Pt(2.2)
    add_textbox(s, 1.05, 4.25, 24.0, 0.7, "Data Pipeline Blueprint", size=30, color="teal", bold=True)
    add_textbox(s, 1.1, 5.05, 25.0, 0.55, "Due scenari TO BE per una blueprint riusabile su Kiron CDG, CDG interno e ING ProSIGNAL", size=14, color="navy")
    add_textbox(s, 1.1, 6.6, 14.0, 0.35, "30 giugno 2026", size=9, color="muted")


def context_slide(prs):
    s = blank_slide(prs)
    add_header(s, "Contesto / Esigenza / Obiettivi", 2)
    img = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(0.8), Cm(1.2), Cm(7.0), Cm(5.55))
    img.fill.solid(); img.fill.fore_color.rgb = COLORS["light"]; img.line.color.rgb = COLORS["line"]
    add_textbox(s, 1.25, 1.65, 6.1, 0.45, "Blueprint riusabile", size=17, color="teal", bold=True)
    add_multiline(s, 1.25, 2.45, 5.9, 3.8, [
        "Decisione di piattaforma, non di singolo progetto.",
        "Priorita': ProSIGNAL e Kiron.",
        "Confronto neutrale tra due scenari."
    ], size=9.5, color="navy")
    items = [
        ("1", "Contesto", "Tre iniziative data-intensive richiedono ingestion, controlli, trasformazioni, mart e dashboard governati."),
        ("2", "Esigenza", "Evitare che ogni progetto ricostruisca pipeline, mapping, controlli e output con logiche diverse."),
        ("3", "Obiettivi", "Discutere una blueprint riusabile, confrontando neutralmente Dagster/dbt/Metabase e Talend/Qlik.")
    ]
    y = 1.3
    for n, title, body in items:
        circ = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Cm(9.1), Cm(y + 0.05), Cm(0.65), Cm(0.65))
        circ.fill.solid(); circ.fill.fore_color.rgb = COLORS["teal"]; circ.line.fill.background()
        add_textbox(s, 9.25, y + 0.2, 0.35, 0.25, n, size=9, color="white", bold=True, align=PP_ALIGN.CENTER)
        add_textbox(s, 10.0, y, 5.2, 0.35, title, size=13, color="navy", bold=True)
        add_textbox(s, 10.0, y + 0.45, 20.6, 0.62, body, size=10, color="text")
        y += 1.65


def pipeline_map(prs):
    s = blank_slide(prs)
    add_header(s, "TO BE | Architettura logica", 3)
    add_title(s, "Una buona pipeline separa le responsabilita' prima di scegliere i tool",
              "La blueprint deve rendere espliciti i passaggi comuni: acquisire, verificare, preparare, modellare, controllare, esporre e governare.")
    steps = [
        ("Sorgenti", "file, API, DB"),
        ("Ingestion", "landing controllata"),
        ("Verifica", "data contracts"),
        ("Preparazione", "normalizzazione"),
        ("Mart", "output data product"),
        ("Controllo", "quality gates"),
        ("Esposizione", "BI, export"),
        ("Governance", "audit, run, accessi"),
    ]
    x0, y = 1.0, 3.0
    for i, (t, b) in enumerate(steps):
        x = x0 + i * 3.9
        rounded_rect(s, x, y, 3.0, 1.2, fill="white", line="line")
        add_textbox(s, x + 0.15, y + 0.18, 2.7, 0.28, t, size=10, color="navy", bold=True, align=PP_ALIGN.CENTER)
        add_textbox(s, x + 0.15, y + 0.62, 2.7, 0.3, b, size=8, color="muted", align=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            connector(s, x + 3.0, y + 0.6, x + 3.65, y + 0.6)
    card(s, 1.0, 5.25, 9.7, 1.1, "Principio comune", "raw/source-faithful, trasformazioni tracciabili e controlli dichiarati prima della dashboard.", accent="teal_dark", body_size=9)
    card(s, 11.2, 5.25, 9.7, 1.1, "Valore per i progetti", "stessa spina dorsale per Kiron, ProSIGNAL e CDG interno, con estensioni specifiche per sorgenti e regole.", accent="blue", body_size=9)
    card(s, 21.4, 5.25, 9.7, 1.1, "Decisione da abilitare", "scegliere il modello operativo piu' adatto: controllo ingegneristico o piattaforma enterprise.", accent="green", body_size=9)


def scenarios_intro(prs):
    s = blank_slide(prs)
    add_header(s, "TO BE | Due scenari", 4)
    add_title(s, "Le due opzioni condividono principi, ma cambiano modello operativo",
              "Il confronto non cerca un vincitore teorico: deve chiarire quale blueprint conviene adottare e con quali vincoli.")
    card(s, 1.0, 2.5, 8.7, 3.2, "Parti comuni", [
        "data contracts e mapping",
        "quality gates e scarti",
        "raw/staging/mart",
        "audit, run history e accessi",
        "dashboard, export e reporting"
    ], accent="teal")
    card(s, 10.4, 2.5, 8.7, 3.2, "Scenario A", [
        "Dagster per orchestrazione",
        "dbt per modelli e test",
        "Metabase per analytics",
        "AWS ECS/EC2, RDS o DB cliente",
        "alto controllo engineering"
    ], accent="blue")
    card(s, 19.8, 2.5, 8.7, 3.2, "Scenario B", [
        "Talend/Qlik per integrazione",
        "Qlik Cloud Analytics Premium",
        "data quality product-oriented",
        "50 GB come baseline economics",
        "maggiore dipendenza da entitlement"
    ], accent="green")
    add_textbox(s, 1.0, 6.25, 28.0, 0.4, "Criteri di confronto: riuso, scalabilita', costi, rapidita', complessita', skill, lock-in, blocker e fit sui tre casi.", size=10, color="navy", bold=True)


def scenario_a(prs, n):
    s = blank_slide(prs)
    add_header(s, "Scenario A | Dagster + dbt + Metabase", n)
    add_title(s, "Lo scenario AWS privilegia controllo, componibilita' e riuso engineering",
              "Ogni componente della pipeline resta esplicita e governabile, con trade-off operativo su setup, run e competenze.")
    lanes = [
        ("Sorgenti", "Jira, Tempo, SAP, Excel, file fixed-column"),
        ("Compute AWS", "ECS o EC2; no EKS"),
        ("Orchestrazione", "Dagster job, asset, schedule, sensor"),
        ("Transform", "dbt models, SQL/Python custom"),
        ("Serving", "mart, Metabase, export")
    ]
    y = 2.4
    for i, (t, b) in enumerate(lanes):
        x = 1.1 + i * 6.0
        rounded_rect(s, x, y, 5.0, 1.35, fill="white", line="line")
        add_textbox(s, x + 0.15, y + 0.18, 4.7, 0.3, t, size=10, color="navy", bold=True, align=PP_ALIGN.CENTER)
        add_textbox(s, x + 0.25, y + 0.68, 4.5, 0.4, b, size=8.2, color="muted", align=PP_ALIGN.CENTER)
        if i < len(lanes) - 1:
            connector(s, x + 5.0, y + 0.68, x + 5.8, y + 0.68, color="teal")
    card(s, 1.1, 4.65, 8.8, 1.45, "Benefici", ["riuso tramite asset e modelli", "controllo su pipeline e deploy", "osservabilita' tecnica dei run"], accent="green", body_size=8.5)
    card(s, 10.5, 4.65, 8.8, 1.45, "Limiti", ["richiede data engineering e AWS", "setup e manutenzione runtime", "Metabase da validare per governance cliente"], accent="amber", body_size=8.5)
    card(s, 19.9, 4.65, 8.8, 1.45, "Blocker da chiudere", ["RDS o DB cliente per progetto", "access control e publishing BI", "stima cloud/run"], accent="red", body_size=8.5)


def scenario_b(prs, n):
    s = blank_slide(prs)
    add_header(s, "Scenario B | Talend + Qlik", n)
    add_title(s, "Lo scenario Qlik concentra integrazione e analytics in una filiera piu' product-oriented",
              "Il punto chiave non e' solo il tool: servono entitlement, limiti di capacita' e runtime compatibili con i casi target.")
    lanes = [
        ("Sorgenti", "file, API, DB, sistemi cliente"),
        ("Talend / Qlik DI", "ingestion, parsing, quality"),
        ("Storage / Mart", "DB cliente o warehouse da definire"),
        ("Qlik Cloud", "analytics, dashboard, reload"),
        ("Governance", "tenant, accessi, monitoraggio")
    ]
    y = 2.4
    for i, (t, b) in enumerate(lanes):
        x = 1.1 + i * 6.0
        rounded_rect(s, x, y, 5.0, 1.35, fill="white", line="line")
        add_textbox(s, x + 0.15, y + 0.18, 4.7, 0.3, t, size=10, color="navy", bold=True, align=PP_ALIGN.CENTER)
        add_textbox(s, x + 0.25, y + 0.68, 4.5, 0.4, b, size=8.2, color="muted", align=PP_ALIGN.CENTER)
        if i < len(lanes) - 1:
            connector(s, x + 5.0, y + 0.68, x + 5.8, y + 0.68, color="teal")
    card(s, 1.1, 4.65, 8.8, 1.45, "Benefici", ["fit enterprise se Qlik e' disponibile", "meno custom se capability coperte", "analytics/presentation forte"], accent="green", body_size=8.5)
    card(s, 10.5, 4.65, 8.8, 1.45, "Limiti", ["entitlement e limiti da confermare", "lock-in piu' alto", "fit ProSIGNAL da stressare"], accent="amber", body_size=8.5)
    card(s, 19.9, 4.65, 8.8, 1.45, "Baseline economics", ["piano 50 GB: 30-50k EUR", "range utente, non listino", "periodo/IVA/bundle da chiarire"], accent="blue", body_size=8.5)


def mapping_slide(prs, n, title, scenario):
    s = blank_slide(prs)
    add_header(s, scenario, n)
    add_title(s, title)
    rows = [
        ("Ingestion", "Dagster job / custom connector", "Talend / Qlik data integration"),
        ("Orchestrazione", "Dagster schedule, sensor, asset graph", "Talend runtime, task/reload API Qlik"),
        ("Trasformazioni", "dbt + SQL/Python dove serve", "Talend transformation o DB/Qlik layer"),
        ("Data quality", "dbt tests, asset checks, controlli custom", "Talend/Qlik quality capability, da verificare"),
        ("Presentation", "Metabase dashboard, query, alert, export", "Qlik Cloud Analytics Premium"),
    ]
    x = [1.0, 8.0, 19.0]
    widths = [6.2, 10.2, 10.2]
    headers = ["Componente", "Scenario A", "Scenario B"]
    for i, h in enumerate(headers):
        band = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(x[i]), Cm(2.25), Cm(widths[i]), Cm(0.55))
        band.fill.solid(); band.fill.fore_color.rgb = COLORS["teal" if i == 0 else "navy"]; band.line.fill.background()
        add_textbox(s, x[i] + 0.12, 2.38, widths[i] - 0.24, 0.22, h, size=8.5, color="white", bold=True, align=PP_ALIGN.CENTER)
    y = 2.9
    for comp, a, b in rows:
        for i, val in enumerate([comp, a, b]):
            rounded_rect(s, x[i], y, widths[i], 0.72, fill="very_light" if i == 0 else "white", line="gray", radius=False)
            add_textbox(s, x[i] + 0.12, y + 0.16, widths[i] - 0.24, 0.28, val, size=8.4, color="navy" if i == 0 else "text", bold=(i == 0))
        y += 0.78
    add_textbox(s, 1.0, 6.75, 28.5, 0.3, "Nota: i dettagli Qlik/Talend su entitlement, capacity e runtime restano da confermare prima di una scelta esecutiva.", size=8.8, color="muted")


def comparison(prs):
    s = blank_slide(prs)
    add_header(s, "Confronto scenari", 9)
    add_title(s, "Il trade-off principale e' tra controllo engineering e piattaforma enterprise")
    criteria = [
        ("Riuso", "Alto con asset/modelli standard", "Alto se template e tenant sono riusabili"),
        ("Scalabilita'", "Dipende da AWS, DB e orchestrazione", "Dipende da capacity e runtime Qlik/Talend"),
        ("Costi", "Cloud/run/competenze", "Subscription 50 GB: 30-50k EUR da confermare"),
        ("Rapidita'", "Alta se si riusa CDG", "Alta se licenze e connettori sono pronti"),
        ("Lock-in", "Medio-basso", "Medio-alto"),
        ("Fit ProSIGNAL", "Forte per processing custom", "Da stressare su file grandi/fixed-column"),
    ]
    x = [1.0, 7.3, 19.2]
    widths = [5.4, 10.7, 10.7]
    for i, h in enumerate(["Criterio", "Dagster / dbt / Metabase", "Talend / Qlik"]):
        band = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(x[i]), Cm(2.05), Cm(widths[i]), Cm(0.5))
        band.fill.solid(); band.fill.fore_color.rgb = COLORS["teal" if i != 1 else "blue"]; band.line.fill.background()
        add_textbox(s, x[i]+0.1, 2.17, widths[i]-0.2, 0.22, h, size=8, color="white", bold=True, align=PP_ALIGN.CENTER)
    y = 2.65
    for row in criteria:
        for i, val in enumerate(row):
            rounded_rect(s, x[i], y, widths[i], 0.62, fill="very_light" if i == 0 else "white", line="gray", radius=False)
            add_textbox(s, x[i]+0.12, y+0.14, widths[i]-0.24, 0.24, val, size=7.8, color="navy" if i == 0 else "text", bold=(i == 0))
        y += 0.68


def projection(prs):
    s = blank_slide(prs)
    add_header(s, "Proiezione sui casi", 10)
    add_title(s, "ProSIGNAL e Kiron guidano la priorita', il CDG interno consolida la reference")
    projects = [
        ("ING ProSIGNAL", "Alta", "Forte per file processing custom e controlli cross-file.", "Forte se capability Talend/Qlik reggono fixed-column e volumi."),
        ("Kiron CDG", "Alta", "Buon fit CDG-like con fonti Campus/Zucchetti e output Actual/Forecast.", "Buon fit se Qlik e' preferito e Talend rende trasparenti regole e mapping."),
        ("CDG interno", "Media", "Molto coerente con la reference implementation corrente.", "Possibile alternativa, ma meno allineata allo stato attuale."),
    ]
    y = 2.0
    for name, prio, a, b in projects:
        card(s, 1.0, y, 8.3, 1.35, name, [f"Priorita': {prio}"], accent="teal", body_size=8.5)
        card(s, 9.8, y, 9.7, 1.35, "Scenario A", [a], accent="blue", body_size=8.5)
        card(s, 20.0, y, 9.7, 1.35, "Scenario B", [b], accent="green", body_size=8.5)
        y += 1.55
    add_textbox(s, 1.0, 6.85, 28.0, 0.25, "Dove mancano esempi o chiarimenti, i contenuti restano ipotesi/best practice e non claim cliente.", size=8.5, color="muted")


def wbs_slide(prs):
    s = blank_slide(prs)
    add_header(s, "WBS | Processi e funzioni", 11)
    add_title(s, "Le WBS condividono una spina dorsale comune, ma cambiano sorgenti e controlli")
    center = (15.8, 3.75)
    rounded_rect(s, center[0]-2.0, center[1]-0.45, 4.0, 0.9, fill="light", line="teal")
    add_textbox(s, center[0]-1.85, center[1]-0.12, 3.7, 0.25, "Blueprint pipeline", size=10, color="navy", bold=True, align=PP_ALIGN.CENTER)
    branches = [
        ("Discovery", 4.0, 1.7), ("Data contracts", 8.3, 1.25), ("Ingestion", 22.5, 1.25), ("Quality gates", 27.0, 1.7),
        ("Transform", 4.0, 5.55), ("Mart/output", 8.3, 6.05), ("Presentation", 22.5, 6.05), ("Run/governance", 27.0, 5.55)
    ]
    for label, x, y in branches:
        connector(s, center[0], center[1], x+1.8, y+0.35, color="line")
        rounded_rect(s, x, y, 3.6, 0.72, fill="white", line="line")
        add_textbox(s, x+0.1, y+0.2, 3.4, 0.22, label, size=8.3, color="navy", bold=True, align=PP_ALIGN.CENTER)
    card(s, 11.0, 5.25, 9.6, 1.25, "Derivazione piano", "analisi -> architettura -> implementazione -> test/parallel run -> go-live -> industrializzazione.", accent="teal_dark", body_size=8.8)


def workplan(prs):
    s = blank_slide(prs)
    add_header(s, "Piano di lavoro", 12)
    add_title(s, "Un piano comune riduce rischio e abilita riuso progressivo")
    phases = [
        ("1", "Discovery", "fonti, vincoli, data contracts"),
        ("2", "Architettura", "scenario, runtime, DB, governance"),
        ("3", "Build", "ingestion, transform, quality, mart"),
        ("4", "Validazione", "parallel run, quadrature, UAT"),
        ("5", "Go-live", "dashboard, export, runbook"),
        ("6", "Riuso", "template, backlog, industrializzazione"),
    ]
    for i, (num, title, desc) in enumerate(phases):
        x = 1.0 + i * 5.1
        pill(s, x, 2.4, 3.7, f"{num}. {title}", fill="light", line="teal")
        add_textbox(s, x, 3.05, 3.7, 0.65, desc, size=8.2, color="text", align=PP_ALIGN.CENTER)
        if i < len(phases)-1:
            connector(s, x+3.7, 2.65, x+4.8, 2.65, color="line")
    card(s, 1.0, 5.05, 8.8, 1.2, "Priorita'", ["ProSIGNAL e Kiron prima del CDG interno"], accent="green", body_size=8.8)
    card(s, 10.4, 5.05, 8.8, 1.2, "Metodo", ["pilot comparativo o scenario selezionato"], accent="blue", body_size=8.8)
    card(s, 19.8, 5.05, 8.8, 1.2, "Output", ["blueprint riusabile e piano di adozione"], accent="teal", body_size=8.8)


def economics(prs):
    s = blank_slide(prs)
    add_header(s, "Economics", 13)
    add_title(s, "Per la blueprint servono range e driver, non una quotazione definitiva")
    card(s, 1.0, 2.2, 8.6, 2.4, "Scenario A", [
        "licenze potenzialmente contenute",
        "costo cloud/run da stimare",
        "maggiore effort engineering",
        "riuso forte se standardizzato"
    ], accent="blue", body_size=8.7)
    card(s, 10.2, 2.2, 8.6, 2.4, "Scenario B", [
        "piano 50 GB: 30-50k EUR",
        "range utente, non listino",
        "capacity, bundle e IVA da chiarire",
        "meno custom se capability coperte"
    ], accent="green", body_size=8.7)
    card(s, 19.4, 2.2, 8.6, 2.4, "Driver comuni", [
        "setup blueprint",
        "delivery per progetto",
        "run, licenze e manutenzione",
        "riuso su iniziative successive"
    ], accent="teal", body_size=8.7)
    rounded_rect(s, 1.0, 5.4, 27.0, 0.95, fill="light", line="teal")
    add_textbox(s, 1.4, 5.65, 26.2, 0.32, "Messaggio: discutere fasce e dipendenze; evitare falsa precisione prima della verifica commerciale e tecnica.", size=9.2, color="navy", bold=True, align=PP_ALIGN.CENTER)


def decision(prs):
    s = blank_slide(prs)
    add_header(s, "Discussione", 14)
    add_title(s, "La scelta deve nascere da una discussione guidata sui trade-off")
    qs = [
        ("Adozione", "Quale scenario puo' diventare standard per progetti analoghi?"),
        ("Rischio", "Quali blocker tecnici o commerciali vanno chiusi prima?"),
        ("Economics", "Il range Qlik/Talend e il costo run AWS sono confrontabili?"),
        ("Priorita'", "Quale caso pilota massimizza valore e apprendimento?")
    ]
    for i, (t, b) in enumerate(qs):
        x = 1.0 + (i % 2) * 14.0
        y = 2.25 + (i // 2) * 1.75
        card(s, x, y, 12.5, 1.25, t, b, accent="teal" if i % 2 == 0 else "blue", body_size=9.5)
    add_textbox(s, 1.0, 6.3, 27.0, 0.4, "La presentazione non chiude la decisione: prepara una scelta informata tra due modelli riusabili.", size=11, color="navy", bold=True, align=PP_ALIGN.CENTER)


def closing(prs):
    s = blank_slide(prs)
    bg = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(0), Cm(0), Cm(33.87), Cm(19.05))
    bg.fill.solid(); bg.fill.fore_color.rgb = COLORS["very_light"]; bg.line.fill.background()
    s.shapes.add_picture(str(ASSET_LOGO), Cm(12.2), Cm(2.65), width=Cm(8.6))
    add_textbox(s, 9.0, 4.4, 15.8, 0.6, "Data Pipeline Blueprint", size=18, color="navy", bold=True, align=PP_ALIGN.CENTER)
    add_textbox(s, 8.2, 5.15, 17.4, 0.5, "Materiale di discussione per CEO, CTO e Tech Committee", size=11, color="muted", align=PP_ALIGN.CENTER)
    add_textbox(s, 1.0, 6.95, 31.0, 0.25, "Documento di lavoro: assunzioni, economics e capability vendor da confermare prima di una proposta definitiva.", size=7.5, color="muted", align=PP_ALIGN.CENTER)


def build():
    ensure_logo()
    prs = Presentation()
    prs.slide_width = Cm(W * 2.54)
    prs.slide_height = Cm(H * 2.54)

    cover(prs)
    context_slide(prs)
    pipeline_map(prs)
    scenarios_intro(prs)
    scenario_a(prs, 5)
    mapping_slide(prs, 6, "Nello scenario AWS ogni componente resta esplicita e governabile", "Scenario A | Mapping componenti")
    scenario_b(prs, 7)
    mapping_slide(prs, 8, "Nello scenario Qlik/Talend capability e analytics sono piu' concentrate", "Scenario B | Mapping componenti")
    comparison(prs)
    projection(prs)
    wbs_slide(prs)
    workplan(prs)
    economics(prs)
    decision(prs)
    closing(prs)

    if OUT.exists():
        raise FileExistsError(f"Output already exists: {OUT}")
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
